"""
core/topology/export/pptx_exporter.py
======================================
Exports a TopologyGraph as a fully editable PowerPoint diagram --
every device is a real movable shape, every link a real connector,
not a flattened image. Opens in PowerPoint/Keynote/Google Slides with
all shapes selectable and re-positionable.

Slide size is computed per-topology via recommended_canvas_size() --
a small lab diagram gets a normal 13.33x7.5 slide, while a 70+ device
site gets a proportionally larger slide (capped at 50", safely under
PowerPoint's hard 56" limit) so nodes never overlap at scale.
"""
from __future__ import annotations

import io
import logging
from typing import Optional

from core.topology.topology_models import TopologyGraph, DeviceRole
from core.topology.export.coords import compute_canvas_positions
from core.topology.layout import (
    recommended_canvas_size, elbow_path, compute_link_slot_fractions, compute_bend_fractions,
)
from core.topology.interface_naming import abbreviate_interface

logger = logging.getLogger("NetBrain.Topology.Export.PPTX")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False


NODE_W_IN  = 1.7
NODE_H_IN  = 0.65


def _hex_to_rgbcolor(hex_str: str) -> "RGBColor":
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def export_topology_to_pptx(graph: TopologyGraph) -> Optional[bytes]:
    """
    Build a .pptx file in memory and return its bytes, or None if the
    python-pptx library isn't available.
    """
    if not PPTX_OK:
        logger.warning("python-pptx not installed -- PPTX export unavailable")
        return None
    if not graph.nodes:
        return None

    slide_w_in, slide_h_in = recommended_canvas_size(graph)

    positions = compute_canvas_positions(
        graph,
        canvas_width_in=slide_w_in,
        canvas_height_in=slide_h_in - 0.9,   # leave room for title
        margin_in=0.9,
    )

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # -- Title --
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(slide_w_in - 0.8), Inches(0.6))
    tf = title_box.text_frame
    tf.text = f"Network Topology — {graph.site_name} ({graph.city}, {graph.country})"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)

    y_offset = 0.9  # push everything below the title

    # -- Links (draw first so they sit behind node shapes) --
    # Orthogonal elbow routing + per-endpoint port labels, slotted per
    # device so multiple links sharing one device don't have their
    # labels collide -- same convention used by the interactive view
    # and the PDF/Visio exports, matching standard tree-style Visio
    # network diagrams.
    slot_fractions = compute_link_slot_fractions(graph)
    bend_fractions = compute_bend_fractions(graph)
    for idx, link in enumerate(graph.links):
        if link.device_a_ip not in positions or link.device_b_ip not in positions:
            continue
        ax, ay = positions[link.device_a_ip]
        bx, by = positions[link.device_b_ip]

        a_frac, b_frac = slot_fractions.get(idx, (0.0, 0.0))
        ax_anchor = ax + NODE_W_IN / 2 + a_frac * NODE_W_IN
        bx_anchor = bx + NODE_W_IN / 2 + b_frac * NODE_W_IN
        ay_anchor = ay + y_offset + NODE_H_IN / 2
        by_anchor = by + y_offset + NODE_H_IN / 2

        path = elbow_path(ax_anchor, ay_anchor, bx_anchor, by_anchor, bend_fraction=bend_fractions.get(idx, 0.5))
        fb = slide.shapes.build_freeform(start_x=path[0][0], start_y=path[0][1], scale=914400)
        fb.add_line_segments(path[1:], close=False)
        conn = fb.convert_to_shape()
        conn.line.color.rgb = RGBColor(0x64, 0x74, 0x8b)
        conn.line.width = Pt(1.5)
        conn.fill.background()

        # Per-endpoint port labels, positioned a short distance along
        # each end's own segment of the path -- not a single combined
        # label at the midpoint.
        a_lx = path[0][0] + (path[1][0] - path[0][0]) * 0.3
        a_ly = path[0][1] + (path[1][1] - path[0][1]) * 0.3
        b_lx = path[-1][0] + (path[-2][0] - path[-1][0]) * 0.3
        b_ly = path[-1][1] + (path[-2][1] - path[-1][1]) * 0.3

        for lx, ly, port in (
            (a_lx, a_ly, link.device_a_port),
            (b_lx, b_ly, link.device_b_port),
        ):
            label_tb = slide.shapes.add_textbox(
                Inches(lx - 0.4), Inches(ly - 0.1), Inches(0.8), Inches(0.2)
            )
            label_tf = label_tb.text_frame
            label_tf.text = abbreviate_interface(port)
            label_tf.paragraphs[0].font.size = Pt(7)
            label_tf.paragraphs[0].font.color.rgb = RGBColor(0x47, 0x55, 0x69)
            label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # -- Nodes --
    for ip, node in graph.nodes.items():
        if ip not in positions:
            continue
        px, py = positions[ip]
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px), Inches(py + y_offset),
            Inches(NODE_W_IN), Inches(NODE_H_IN),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex_to_rgbcolor(node.role.color)
        shp.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shp.line.width = Pt(1)

        tf = shp.text_frame
        tf.word_wrap = True
        tf.text = f"{node.role.icon} {node.label()}"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = ip if not ip.startswith("unknown:") else "(no IP reported)"
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        p2.alignment = PP_ALIGN.CENTER

    # -- Legend --
    legend_y = slide_h_in - 0.5
    legend_x = 0.4
    for role in (DeviceRole.ROUTER, DeviceRole.SWITCH, DeviceRole.ACCESS_POINT, DeviceRole.FIREWALL):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(legend_x), Inches(legend_y), Inches(0.18), Inches(0.18)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _hex_to_rgbcolor(role.color)
        dot.line.fill.background()

        lbl = slide.shapes.add_textbox(Inches(legend_x + 0.22), Inches(legend_y - 0.05), Inches(1.5), Inches(0.3))
        lbl.text_frame.text = role.value.replace("_", " ").title()
        lbl.text_frame.paragraphs[0].font.size = Pt(10)
        legend_x += 1.8

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
